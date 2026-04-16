#!/usr/bin/env python3
"""
Generate Business Model PowerPoint Presentation for Smart Green Space
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Color scheme (green theme for green space)
COLOR_PRIMARY = RGBColor(34, 139, 34)      # Forest Green
COLOR_SECONDARY = RGBColor(50, 205, 50)   # Lime Green
COLOR_ACCENT = RGBColor(0, 100, 0)        # Dark Green
COLOR_TEXT = RGBColor(40, 40, 40)         # Dark Gray
COLOR_LIGHT = RGBColor(240, 248, 240)     # Ghost White

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_LIGHT
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.5), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    p = left_title_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(4.8))
    left_frame = left_content.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = "• " + item
        p.font.size = Pt(14)
        p.space_before = Pt(8)
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    p = right_title_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    right_content = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(4.8))
    right_frame = right_content.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = "• " + item
        p.font.size = Pt(14)
        p.space_before = Pt(8)

def create_presentation():
    """Create the entire presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
                   "Smart Green Space",
                   "Transform Urban Resilience Through Real-Time Green & Thermal Intelligence")
    
    # Slide 2: The Problem
    add_content_slide(prs,
                     "The Challenge",
                     [
                         "🌍 Urban flooding & heat stress threaten 80% of global cities",
                         "📊 City planners lack real-time insights on green space health",
                         "⚠️ Emergency response remains reactive, not proactive",
                         "💰 $3T annual infrastructure loss from climate-related hazards",
                         "🖥️ Fragmented data sources = no single source of truth"
                     ])
    
    # Slide 3: Our Solution
    add_content_slide(prs,
                     "Our Solution",
                     [
                         "✅ Proactive Resilience Platform transforming civic action from reactive to proactive",
                         "✅ GSHI Score: Green Space Health Index using 7-component AI model",
                         "✅ Real-time satellite + sensor data aggregation (satellite, weather, IoT)",
                         "✅ 3D visual interface making complex risk intuitive & actionable",
                         "✅ API-first architecture for seamless integration"
                     ])
    
    # Slide 4: Key Partners
    add_content_slide(prs,
                     "Key Partners & Resources",
                     [
                         "📡 Data Providers: Open-Meteo, TomTom, Mapbox, Google Earth Engine",
                         "🏛️ Civic Bodies: Municipal corporations, meteorological departments",
                         "☁️ Cloud Providers: AWS, DigitalOcean (Backend), Vercel (Edge UI)",
                         "🤖 Proprietary ML: Random forest architectures for flood prediction",
                         "👥 Human Capital: GIS specialists, DevOps engineers, ML researchers"
                     ])
    
    # Slide 5: Value Propositions
    add_two_column_slide(prs,
                        "Value Propositions",
                        "For Government",
                        [
                            "Proactive disaster navigation",
                            "Evidence-based urban planning",
                            "Real-time emergency preparedness",
                            "Cost optimization & prevention"
                        ],
                        "For Enterprise",
                        [
                            "Zone risk assessment for development",
                            "Flood liability quantification for insurance",
                            "Long-term investment protection",
                            "Regulatory compliance support"
                        ])
    
    # Slide 6: Business Model
    add_content_slide(prs,
                     "Revenue Streams",
                     [
                         "💼 SaaS Subscriptions: Tiered monthly access ($500 - $5k/mo)",
                         "📜 Civic Contracts: Multi-year licensing for city digital twins",
                         "🔗 API Tokens: Metered billing for underwriter algorithms",
                         "🔧 Custom Integrations: Premium deployment & support services"
                     ])
    
    # Slide 7: Customer Segments
    add_two_column_slide(prs,
                        "Customer Base",
                        "B2G Segment",
                        [
                            "City planning departments",
                            "Emergency response hubs (NDRF)",
                            "Municipal corporations",
                            "Climate adaptation committees"
                        ],
                        "B2B Segment",
                        [
                            "Property developers",
                            "Insurance companies",
                            "Climate tech investors",
                            "Sustainability consultants"
                        ])
    
    # Slide 8: Key Activities
    add_content_slide(prs,
                     "Key Activities",
                     [
                         "🔧 Platform Development: Continuous ML model training & optimization",
                         "📊 Data Aggregation: ETL pipelines for satellite & IoT telemetry",
                         "🎨 UI/UX Rendering: 3D visualization & map element optimization",
                         "📈 Model Training: Random forest algorithms for health & flood prediction",
                         "🚀 Deployment: Continuous SaaS updates & new geographical zones"
                     ])
    
    # Slide 9: Cost Structure
    add_content_slide(prs,
                     "Cost Structure",
                     [
                         "⚙️ Infrastructure: GPU clusters, PostgreSQL hosting, cloud computing",
                         "📜 Licensing: Premium mapping tiles, real-time API quotas",
                         "👨‍💻 R&D & Payroll: Algorithm development, frontend iteration",
                         "📡 Data: Satellite imagery & weather data subscriptions",
                         "🔐 Security & Compliance: Enterprise-grade infrastructure"
                     ])
    
    # Slide 10: Channels
    add_content_slide(prs,
                     "Go-to-Market Channels",
                     [
                         "🎯 Direct Procurement: Government tenders & smart-city grants",
                         "🔗 API Syndication: Developer portal for risk algorithms",
                         "🤝 Partnership: Insurance brokers, real estate platforms",
                         "📱 SaaS Dashboard: Direct B2B subscriptions",
                         "🌐 Cloud Marketplaces: AWS, Google Cloud integrations"
                     ])
    
    # Slide 11: Customer Relationships
    add_content_slide(prs,
                     "Customer Relationships",
                     [
                         "👥 Dedicated Support: Direct integration teams for enterprise deployment",
                         "📬 Automated Updates: Continuous SaaS improvements & new ML models",
                         "📞 Concierge Service: Custom onboarding & training programs",
                         "🔔 Real-time Alerts: Event-driven notifications & anomaly detection",
                         "📊 Analytics Dashboard: Usage insights & performance metrics"
                     ])
    
    # Slide 12: GSHI Components
    add_content_slide(prs,
                     "Green Space Health Index (GSHI) - 7 Components",
                     [
                         "🌿 Vegetation (22%): Satellite NDVI normalized for urban canopy",
                         "🌡️ Thermal (18%): Heat index from temperature + humidity",
                         "💧 Water (17%): Soil moisture & drainage health",
                         "🦅 Biodiversity (15%): Shannon diversity with conservation weighting",
                         "💨 Air Quality (10%): PM2.5, PM10, CO2 concentrations",
                         "⚡ Infrastructure (9%): Sensor uptime & alert responsiveness",
                         "🌲 Tree Health (9%): AI-scanned decay analysis"
                     ])
    
    # Slide 13: Technology Stack
    add_content_slide(prs,
                     "Technology Stack",
                     [
                         "🎯 Frontend: React/TypeScript with 3D visualization (Three.js, Mapbox)",
                         "🔧 Backend: Node.js with Express & Prisma ORM",
                         "📊 Database: PostgreSQL (TimescaleDB for time-series)",
                         "🤖 ML: Python (scikit-learn, TensorFlow) on GPU clusters",
                         "☁️ Cloud: AWS EC2, Docker containerization, Redis caching"
                     ])
    
    # Slide 14: Traction & Roadmap
    add_content_slide(prs,
                     "Roadmap & Milestones",
                     [
                         "Q2 2026: Alpha launch for 3 cities (Delhi, Mumbai, Bangalore)",
                         "Q3 2026: Beta SaaS platform & API documentation",
                         "Q4 2026: Enterprise sales & first civic contracts",
                         "2027: 10+ city deployments & insurance partnership pilots",
                         "2028: Pan-India expansion & international licensing deals"
                     ])
    
    # Slide 15: Call to Action
    add_title_slide(prs,
                   "Building Climate-Resilient Cities",
                   "Partner with us to transform urban sustainability")
    
    # Save presentation
    output_path = r"c:\Users\ayush\OneDrive\Documents\algo trading (2)\smart-green-space\Business_Model_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
